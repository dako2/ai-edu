"""
#from services.google_slides import *
from services.pptx_slides import *

credentials_file = "secrets/closeby-440718-dd98e45706c2.json"
presentation_id = "1cnaNzoYKHz2A-gbJ-9EJEaQT8L-uN2FDzNQOBy3lIbE"

#slide_service = GoogleSlideService(credentials_file, presentation_id)
slide_service = PptxSlideService("ppts/鞠化学发展史-第一章.pptx")

#slides_obj = SlideCollection.load_from_json("services/assets/pptx_slides_temp.json")
slides_obj = SlideCollection.load_from_json("鞠化学发展史-第一章.pptx.json")
slides_obj.tts()
slides_obj.save_to_json("鞠化学发展史-第一章_tts.json")
"""

import json
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE


# Load JSON data
with open('services/assets/鞠化学发展史-第一章_tts.json', 'r', encoding='utf-8') as file:
    audio_data = json.load(file)

# Load the presentation
presentation_path = 'ppts/鞠化学发展史-第一章.pptx'
presentation = Presentation(presentation_path)

for entry in audio_data:
    slide_index = entry['slide_index']
    audio_file = entry['audio_file']
    is_generated = entry['is_generated']

    # Check if audio file exists and is generated
    if audio_file and is_generated:
        # Ensure the slide index is within the valid range
        if 0 <= slide_index < len(presentation.slides):
            slide = presentation.slides[slide_index]

            # Add a shape to represent the audio icon
            left = Inches(1)  # Adjust position as needed
            top = Inches(1)
            width = Inches(1)
            height = Inches(1)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )

            # Add a hyperlink to the audio file
            shape.click_action.hyperlink.address = audio_file

            # Optionally, set the shape's text or appearance
            shape.text = "Play Audio"

# Save the updated presentation
updated_presentation_path = 'ppts/updated_xx.pptx'
presentation.save(updated_presentation_path)
