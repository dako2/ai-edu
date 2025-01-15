# pptx_service.py
from pptx import Presentation
from core.slide_data import SlideData, SlideCollection  # Assuming you have a unified data structure

class PptxSlideService:
    def __init__(self, pptx_file):
        self.presentation = Presentation(pptx_file)
        # Fetch slide data as a list of SlideData instances
        self.presentation_id = pptx_file.split("/")[-1]
        self.slides = self.get_slides()
        

    def get_slides(self, presentation_id=None):
        if not presentation_id:
            presentation_id = self.presentation_id
        else:
            self.presentation_id = presentation_id.split("/")[-1]

        """Fetch all slides in the presentation, including speaker notes and slide IDs."""
        self.slide_data_list = SlideCollection()
        for idx, slide in enumerate(self.presentation.slides):
            speaker_notes = self._get_speaker_notes(slide)
            slide_obj = SlideData(
                slide_index=idx,
                presentation_id=self.presentation_id,
                object_id=str(idx),  # Use the slide index as a unique identifier
                speaker_notes=speaker_notes
            )
            self.slide_data_list.add_slide(slide_obj)
        self.slide_data_list.save_to_json(self.presentation_id+".json")
        return self.slide_data_list

    def _get_speaker_notes(self, slide):
        """Extract speaker notes from a slide."""
        try:
            notes_slide = slide.notes_slide
        except AttributeError:
            # Slide might not have a notes_slide property
            return ""
        
        text = ""
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                text += shape.text
        return text

    def add_slide(self, layout_index=0, speaker_notes=""):
        """
        Add a new slide with optional speaker notes.
        layout_index: Index of the slide layout to use.
        """
        # Select a slide layout (default to first layout if not specified)
        slide_layout = self.presentation.slide_layouts[layout_index]
        new_slide = self.presentation.slides.add_slide(slide_layout)

        # Add speaker notes to the new slide if provided
        if speaker_notes:
            notes_slide = new_slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = speaker_notes

        # Update internal slide list
        self.slides = self.get_slides()
        self.total_slides = len(self.slides)
        return new_slide

    def save(self, output_file):
        """Save the current presentation to a file."""
        self.presentation.save(output_file)

    def show_slide_url(self, slide_index):
        """
        Since PPTX is a local file, you can't generate a URL like Google Slides.
        Instead, this could return a file path or a message indicating the slide is loaded.
        """
        if 0 <= slide_index < len(self.slides):
            return f"Slide {slide_index + 1} loaded from PPTX."
        else:
            print(f"Invalid slide index: {slide_index}.")
            return None

    def get_speaker_notes_for_slide(self, slide_index):
        """Retrieve the speaker notes for a specific slide."""
        if 0 <= slide_index < len(self.slides):
            return self.slides[slide_index].speaker_notes
        else:
            return {"error": "Invalid slide index."}

# Example usage:
if __name__ == "__main__":
    pptx_file = "/Users/dako22/Downloads/test_full.pptx"
    pptx_service = PptxSlideService(pptx_file)
    print(pptx_service.slides[1].speaker_notes)
